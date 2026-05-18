#!/usr/bin/env python3
"""
Component Loader for Infographic Templates
Loads pre-designed infographic components from .pptx and .odp files.
Supports cloning components into presentations and filling with dynamic content.
"""

import os
import subprocess
import tempfile
from pathlib import Path
from typing import Optional, List, Dict, Any
from copy import deepcopy

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    raise ImportError("python-pptx is required. Install with: pip install python-pptx")


class ComponentLibrary:
    """Manages infographic components from .pptx and .odp files"""

    def __init__(self, components_dir: str = "references/components"):
        """
        Initialize component library.

        Args:
            components_dir: Directory containing .pptx and .odp component files
        """
        self.components_dir = Path(components_dir)
        self.components: Dict[str, Path] = {}
        self._discover_components()

    def _discover_components(self):
        """Discover all .pptx and .odp files in components directory"""
        if not self.components_dir.exists():
            return

        # Find all .pptx files
        for pptx_file in self.components_dir.glob("*.pptx"):
            component_name = pptx_file.stem  # filename without extension
            self.components[component_name] = pptx_file

        # Find all .odp files and convert them
        for odp_file in self.components_dir.glob("*.odp"):
            component_name = odp_file.stem
            pptx_path = self._convert_odp_to_pptx(odp_file)
            if pptx_path:
                self.components[component_name] = pptx_path

    def _convert_odp_to_pptx(self, odp_path: Path) -> Optional[Path]:
        """
        Convert ODP file to PPTX using LibreOffice.

        Args:
            odp_path: Path to .odp file

        Returns:
            Path to converted .pptx file, or None if conversion failed
        """
        try:
            # Create temp directory for converted file
            with tempfile.TemporaryDirectory() as tmpdir:
                output_path = Path(tmpdir) / f"{odp_path.stem}.pptx"

                # Use LibreOffice to convert
                subprocess.run(
                    [
                        "libreoffice",
                        "--headless",
                        "--convert-to",
                        "pptx",
                        "--outdir",
                        tmpdir,
                        str(odp_path),
                    ],
                    check=True,
                    capture_output=True,
                    timeout=30,
                )

                # Copy to components cache
                cache_path = self.components_dir / f"{odp_path.stem}.pptx"
                if output_path.exists():
                    import shutil
                    shutil.copy(output_path, cache_path)
                    return cache_path

        except (subprocess.CalledProcessError, FileNotFoundError, subprocess.TimeoutExpired):
            # LibreOffice not available or conversion failed
            # Return None and skip this component
            pass

        return None

    def list_components(self) -> List[str]:
        """Get list of available component names"""
        return sorted(self.components.keys())

    def get_component(self, component_name: str) -> Optional[Presentation]:
        """
        Load a component presentation.

        Args:
            component_name: Name of the component (without extension)

        Returns:
            Presentation object or None if not found
        """
        if component_name not in self.components:
            return None

        try:
            return Presentation(str(self.components[component_name]))
        except Exception as e:
            print(f"Error loading component '{component_name}': {e}")
            return None

    def clone_component_slide(
        self,
        target_presentation: Presentation,
        component_name: str,
        slide_index: int = 0,
    ) -> Optional[Any]:
        """
        Clone a component slide into a target presentation.

        Args:
            target_presentation: Presentation to add slide to
            component_name: Name of component to clone
            slide_index: Which slide from component to use (default: 0)

        Returns:
            The cloned slide in target presentation, or None if failed
        """
        component = self.get_component(component_name)
        if not component or len(component.slides) <= slide_index:
            return None

        try:
            # Get the component slide
            component_slide = component.slides[slide_index]

            # Add a blank slide to target
            blank_slide_layout = target_presentation.slide_layouts[6]  # Blank layout
            new_slide = target_presentation.slides.add_slide(blank_slide_layout)

            # Copy all shapes from component to new slide
            for shape in component_slide.shapes:
                el = shape.element
                newel = deepcopy(el)
                new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

            return new_slide

        except Exception as e:
            print(f"Error cloning component '{component_name}': {e}")
            return None

    def get_component_text_boxes(self, slide: Any) -> List[Any]:
        """
        Get all text boxes in a slide (for updating content).

        Args:
            slide: Slide object

        Returns:
            List of text box shapes
        """
        text_boxes = []
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                text_boxes.append(shape)
        return text_boxes


class ComponentFiller:
    """Fills component templates with dynamic content"""

    def __init__(self, slide: Any):
        """
        Initialize filler for a cloned component slide.

        Args:
            slide: The slide to fill
        """
        self.slide = slide
        self.text_boxes = self._get_text_boxes()

    def _get_text_boxes(self) -> List[Any]:
        """Extract all text boxes from slide"""
        boxes = []
        for shape in self.slide.shapes:
            if hasattr(shape, "text_frame"):
                boxes.append(shape)
        return boxes

    def fill_metrics(self, metrics: List[Dict[str, Any]]):
        """
        Fill metric values in metric component.

        Args:
            metrics: List of dicts with 'key', 'value', 'change' fields
        """
        for i, metric in enumerate(metrics):
            if i >= len(self.text_boxes):
                break

            text_box = self.text_boxes[i]
            value = metric.get("value", "")
            change = metric.get("change", "")

            # Format: "Value\nChange"
            text_box.text_frame.clear()
            p = text_box.text_frame.paragraphs[0]
            p.text = value

            if change:
                p_change = text_box.text_frame.add_paragraph()
                p_change.text = change
                p_change.level = 0

    def fill_features(self, features: List[Dict[str, str]]):
        """
        Fill feature items in feature component.

        Args:
            features: List of dicts with 'title', 'description' fields
        """
        for i, feature in enumerate(features):
            if i >= len(self.text_boxes):
                break

            text_box = self.text_boxes[i]
            title = feature.get("title", "")
            description = feature.get("description", "")

            # Format: "Title\nDescription"
            text_box.text_frame.clear()
            p = text_box.text_frame.paragraphs[0]
            p.text = title

            if description:
                p_desc = text_box.text_frame.add_paragraph()
                p_desc.text = description
                p_desc.level = 1

    def fill_comparison(self, left_items: List[str], right_items: List[str]):
        """
        Fill left/right comparison items.

        Args:
            left_items: Items for left column
            right_items: Items for right column
        """
        # Assuming text boxes are ordered: left items, then right items
        mid = len(self.text_boxes) // 2

        for i, item in enumerate(left_items):
            if i < mid:
                self.text_boxes[i].text_frame.text = item

        for i, item in enumerate(right_items):
            if i + mid < len(self.text_boxes):
                self.text_boxes[i + mid].text_frame.text = item

    def fill_impact_statement(self, statement: str, subtitle: str = ""):
        """
        Fill impact/closing statement.

        Args:
            statement: Main statement text
            subtitle: Optional subtitle
        """
        if self.text_boxes:
            self.text_boxes[0].text_frame.text = statement

        if subtitle and len(self.text_boxes) > 1:
            self.text_boxes[1].text_frame.text = subtitle

    def update_text(self, box_index: int, text: str):
        """Update a specific text box by index"""
        if 0 <= box_index < len(self.text_boxes):
            self.text_boxes[box_index].text_frame.text = text


def create_component_library(components_dir: str = "references/components") -> ComponentLibrary:
    """Factory for creating a ComponentLibrary"""
    return ComponentLibrary(components_dir)
