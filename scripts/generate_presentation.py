#!/usr/bin/env python3
"""
Main presentation generator orchestrator.
Combines markdown parsing, component loading, and design system application.

Usage:
    from scripts.generate_presentation import generate_presentation_from_markdown
    generate_presentation_from_markdown("content.md", "output.pptx")
"""

from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt

from .markdown_parser import parse_markdown, Slide
from .component_loader import create_component_library, ComponentFiller
from .design_system import create_design_applier


def generate_presentation_from_markdown(
    markdown_path: str,
    output_path: str = "presentation.pptx",
    components_dir: str = "references/components",
    design_md_path: str = "references/DESIGN.md",
) -> bool:
    """
    Generate a professional PowerPoint presentation from a markdown file.

    The workflow:
    1. Parse markdown into slides (detects: metrics, features, comparison, impact, text)
    2. Load component library (.pptx and .odp files)
    3. For each slide:
       - Select matching component
       - Clone into main presentation
       - Fill with content
       - Apply design system styling
    4. Save as .pptx

    Args:
        markdown_path: Path to input .md file
        output_path: Path to output .pptx file
        components_dir: Directory containing .pptx/.odp component files
        design_md_path: Path to DESIGN.md for styling rules

    Returns:
        True if successful, False otherwise
    """
    try:
        print(f"📖 Parsing markdown: {markdown_path}")
        title, subtitle, slides = parse_markdown(markdown_path)

        print(f"🎨 Loading components from: {components_dir}")
        library = create_component_library(components_dir)
        available = library.list_components()
        print(f"   Found {len(available)} components: {', '.join(available)}")

        print(f"🎭 Loading design system: {design_md_path}")
        designer = create_design_applier(design_md_path)

        print(f"\n📝 Creating presentation: {title}")
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # Title slide
        if title:
            title_slide_layout = prs.slide_layouts[0]  # Title slide layout
            slide = prs.slides.add_slide(title_slide_layout)
            title_shape = slide.shapes.title
            subtitle_shape = slide.placeholders[1]

            title_shape.text = title
            if subtitle:
                subtitle_shape.text = subtitle

            print(f"  ✓ Title slide: {title}")

        # Content slides
        for i, slide_info in enumerate(slides, 1):
            print(f"  Slide {i}: {slide_info.title} ({slide_info.slide_type})")

            # Map slide type to component
            component_name = _get_component_for_slide(slide_info)

            if not component_name:
                print(f"    ⚠ No component for type '{slide_info.slide_type}', using text slide")
                _add_text_slide(prs, slide_info, designer)
                continue

            # Try to load and use component
            cloned_slide = library.clone_component_slide(prs, component_name)

            if not cloned_slide:
                print(f"    ⚠ Component '{component_name}' not found, using text slide")
                _add_text_slide(prs, slide_info, designer)
                continue

            # Fill the component with content
            filler = ComponentFiller(cloned_slide)
            _fill_slide_content(filler, slide_info)

            print(f"    ✓ Cloned component: {component_name}")

        print(f"\n💾 Saving presentation: {output_path}")
        prs.save(output_path)
        print(f"✨ Done! Presentation saved to {output_path}")

        return True

    except Exception as e:
        print(f"❌ Error: {e}")
        return False


def _get_component_for_slide(slide: Slide) -> Optional[str]:
    """Map slide type to component name"""
    mapping = {
        "metrics": "metric-grid-4",  # Default to 4-metric grid
        "features": "feature-grid-3",  # Default to 3-feature grid
        "comparison": "comparison-panel",
        "impact": "impact-statement",
        "text": None,  # No component, use simple text slide
    }
    return mapping.get(slide.slide_type)


def _fill_slide_content(filler: ComponentFiller, slide: Slide):
    """Fill component with slide content"""
    if slide.slide_type == "metrics":
        # Convert Metric objects to dicts
        metrics_dicts = [
            {
                "key": m.key,
                "value": m.value,
                "change": m.change,
            }
            for m in slide.content
        ]
        filler.fill_metrics(metrics_dicts)

    elif slide.slide_type == "features":
        filler.fill_features(slide.content)

    elif slide.slide_type == "comparison":
        left_items = slide.content.get("left", {}).get("items", [])
        right_items = slide.content.get("right", {}).get("items", [])
        filler.fill_comparison(left_items, right_items)

    elif slide.slide_type == "impact":
        filler.fill_impact_statement(slide.content)

    elif slide.slide_type == "text":
        if isinstance(slide.content, list):
            text = "\n".join(slide.content)
        else:
            text = str(slide.content)
        filler.fill_impact_statement(text)


def _add_text_slide(prs: Presentation, slide: Slide, designer):
    """Add a simple text slide without a component"""
    text_slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide_obj = prs.slides.add_slide(text_slide_layout)

    title_shape = slide_obj.shapes.title
    title_shape.text = slide.title

    if slide.content:
        body_shape = slide_obj.placeholders[1]
        text_frame = body_shape.text_frame
        text_frame.clear()

        if isinstance(slide.content, list):
            for line in slide.content:
                p = text_frame.add_paragraph()
                p.text = line
                p.level = 0
        else:
            text_frame.text = str(slide.content)


if __name__ == "__main__":
    import sys

    if len(sys.argv) < 2:
        print("Usage: python -m scripts.generate_presentation <markdown_file> [output.pptx]")
        sys.exit(1)

    markdown_file = sys.argv[1]
    output_file = sys.argv[2] if len(sys.argv) > 2 else "presentation.pptx"

    success = generate_presentation_from_markdown(markdown_file, output_file)
    sys.exit(0 if success else 1)
