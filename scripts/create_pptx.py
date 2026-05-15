#!/usr/bin/env python3
"""
AmEx PowerPoint Presentation Generator
Generates professional branded PowerPoint presentations with Material Design icons.
"""

import json
import os
from pathlib import Path
from typing import List, Dict, Optional, Tuple
from dataclasses import dataclass
from enum import Enum

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
except ImportError:
    raise ImportError("python-pptx is required. Install with: pip install python-pptx")


class SlideLayout(Enum):
    """Supported slide layouts"""
    TITLE = "title"
    CONTENT = "content"
    TWO_COLUMN = "two_column"
    ICON_GRID = "icon_grid"
    IMAGE_LEFT = "image_left"
    IMAGE_RIGHT = "image_right"
    FULL_IMAGE = "full_image"
    QUOTE = "quote"
    SECTION = "section"


@dataclass
class ColorScheme:
    """Color configuration from color-tokens.json"""
    primary: str
    secondary: str
    accent: str
    dark: str
    light: str
    success: str
    warning: str
    danger: str
    text_primary: str
    text_secondary: str
    bg_primary: str
    bg_secondary: str


class PowerPointGenerator:
    """Main class for generating PowerPoint presentations"""

    def __init__(
        self,
        title: str = "Presentation",
        subtitle: str = "",
        color_tokens_path: Optional[str] = None,
        icons_index_path: Optional[str] = None,
        output_path: str = "presentation.pptx"
    ):
        """
        Initialize the PowerPoint generator.

        Args:
            title: Presentation title
            subtitle: Presentation subtitle
            color_tokens_path: Path to color-tokens.json
            icons_index_path: Path to material-icons-index.json
            output_path: Where to save the .pptx file
        """
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)

        self.title = title
        self.subtitle = subtitle
        self.output_path = output_path

        # Load configuration
        self.colors = self._load_colors(color_tokens_path)
        self.icons_index = self._load_icons_index(icons_index_path)

        # Track slides
        self.slide_count = 0

    def _load_colors(self, path: Optional[str]) -> ColorScheme:
        """Load color tokens from JSON"""
        if not path or not os.path.exists(path):
            # Use defaults if file not found
            return ColorScheme(
                primary="#0066CC",
                secondary="#FF6600",
                accent="#00AA44",
                dark="#1A1A1A",
                light="#F5F5F5",
                success="#00AA44",
                warning="#FFAA00",
                danger="#DD0000",
                text_primary="#1A1A1A",
                text_secondary="#666666",
                bg_primary="#FFFFFF",
                bg_secondary="#F5F5F5"
            )

        with open(path) as f:
            tokens = json.load(f)

        return ColorScheme(
            primary=tokens.get("primary", "#0066CC"),
            secondary=tokens.get("secondary", "#FF6600"),
            accent=tokens.get("accent", "#00AA44"),
            dark=tokens.get("neutral", {}).get("dark", "#1A1A1A"),
            light=tokens.get("neutral", {}).get("light", "#F5F5F5"),
            success=tokens.get("semantic", {}).get("success", "#00AA44"),
            warning=tokens.get("semantic", {}).get("warning", "#FFAA00"),
            danger=tokens.get("semantic", {}).get("danger", "#DD0000"),
            text_primary=tokens.get("text", {}).get("primary", "#1A1A1A"),
            text_secondary=tokens.get("text", {}).get("secondary", "#666666"),
            bg_primary=tokens.get("background", {}).get("primary", "#FFFFFF"),
            bg_secondary=tokens.get("background", {}).get("secondary", "#F5F5F5")
        )

    def _load_icons_index(self, path: Optional[str]) -> Dict:
        """Load Material Design icons index"""
        if not path or not os.path.exists(path):
            return {}

        with open(path) as f:
            return json.load(f)

    def _hex_to_rgb(self, hex_color: str) -> Tuple[int, int, int]:
        """Convert hex color to RGB tuple"""
        hex_color = hex_color.lstrip('#')
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

    def _get_rgb_color(self, hex_color: str) -> RGBColor:
        """Convert hex color to RGBColor"""
        rgb = self._hex_to_rgb(hex_color)
        return RGBColor(*rgb)

    def add_title_slide(self, title: str, subtitle: str = "") -> None:
        """Add a title slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout

        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self._get_rgb_color(self.colors.primary)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(60)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.word_wrap = True
            p = subtitle_frame.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(28)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

        self.slide_count += 1

    def add_content_slide(
        self,
        title: str,
        content: List[str],
        icon: Optional[str] = None,
        layout_type: SlideLayout = SlideLayout.CONTENT
    ) -> None:
        """Add a content slide with bullets"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self._get_rgb_color(self.colors.bg_primary)

        # Header bar with primary color
        header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.2))
        header.fill.solid()
        header.fill.fore_color.rgb = self._get_rgb_color(self.colors.primary)
        header.line.color.rgb = self._get_rgb_color(self.colors.primary)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        # Content area
        content_top = Inches(1.5)
        content_box = slide.shapes.add_textbox(Inches(0.75), content_top, Inches(8.5), Inches(5.5))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        for i, item in enumerate(content):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()

            p.text = item
            p.font.size = Pt(20)
            p.font.color.rgb = self._get_rgb_color(self.colors.text_primary)
            p.level = 0
            p.space_before = Pt(6)
            p.space_after = Pt(6)

        self.slide_count += 1

    def add_two_column_slide(
        self,
        title: str,
        left_title: str,
        left_content: List[str],
        right_title: str,
        right_content: List[str]
    ) -> None:
        """Add a two-column content slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self._get_rgb_color(self.colors.bg_primary)

        # Header bar
        header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.2))
        header.fill.solid()
        header.fill.fore_color.rgb = self._get_rgb_color(self.colors.primary)
        header.line.color.rgb = self._get_rgb_color(self.colors.primary)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        # Left column
        left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(0.4))
        left_title_frame = left_title_box.text_frame
        p = left_title_frame.paragraphs[0]
        p.text = left_title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = self._get_rgb_color(self.colors.primary)

        left_content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(4.5), Inches(4.8))
        left_text_frame = left_content_box.text_frame
        left_text_frame.word_wrap = True

        for i, item in enumerate(left_content):
            if i == 0:
                p = left_text_frame.paragraphs[0]
            else:
                p = left_text_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(16)
            p.font.color.rgb = self._get_rgb_color(self.colors.text_primary)

        # Right column
        right_title_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4), Inches(0.4))
        right_title_frame = right_title_box.text_frame
        p = right_title_frame.paragraphs[0]
        p.text = right_title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = self._get_rgb_color(self.colors.secondary)

        right_content_box = slide.shapes.add_textbox(Inches(5.5), Inches(2.0), Inches(4.5), Inches(4.8))
        right_text_frame = right_content_box.text_frame
        right_text_frame.word_wrap = True

        for i, item in enumerate(right_content):
            if i == 0:
                p = right_text_frame.paragraphs[0]
            else:
                p = right_text_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(16)
            p.font.color.rgb = self._get_rgb_color(self.colors.text_primary)

        self.slide_count += 1

    def add_section_slide(self, title: str, description: str = "") -> None:
        """Add a section divider slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Background - secondary color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self._get_rgb_color(self.colors.secondary)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(60)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # Description
        if description:
            desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
            desc_frame = desc_box.text_frame
            desc_frame.word_wrap = True
            p = desc_frame.paragraphs[0]
            p.text = description
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

        self.slide_count += 1

    def add_closing_slide(self, title: str, cta: str = "") -> None:
        """Add a closing/call-to-action slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self._get_rgb_color(self.colors.primary)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # CTA
        if cta:
            cta_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1))
            cta_frame = cta_box.text_frame
            cta_frame.word_wrap = True
            p = cta_frame.paragraphs[0]
            p.text = cta
            p.font.size = Pt(28)
            p.font.color.rgb = self._get_rgb_color(self.colors.accent)
            p.alignment = PP_ALIGN.CENTER

        self.slide_count += 1

    def save(self) -> str:
        """Save the presentation to file"""
        self.prs.save(self.output_path)
        return self.output_path


def main():
    """Example usage"""
    generator = PowerPointGenerator(
        title="Sample Presentation",
        output_path="sample.pptx"
    )

    generator.add_title_slide("American Express", "Quarterly Review 2024")
    generator.add_section_slide("Executive Summary")
    generator.add_content_slide(
        "Key Metrics",
        [
            "Revenue Growth: +15% YoY",
            "Customer Satisfaction: 94%",
            "Market Share: Increased by 2%"
        ]
    )
    generator.add_two_column_slide(
        "Comparison",
        "2023 Results",
        ["$10M Revenue", "50K Customers", "92% Retention"],
        "2024 Results",
        ["$11.5M Revenue", "52.5K Customers", "94% Retention"]
    )
    generator.add_closing_slide("Thank You", "Questions?")

    output = generator.save()
    print(f"Presentation saved to: {output}")


if __name__ == "__main__":
    main()
